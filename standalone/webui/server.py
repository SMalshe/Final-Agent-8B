"""Agent Lab — a local web console for the per-model agents.

Pick a model, type a task, press Run, and watch the loop work: the plan, each
model call streaming token by token, every tool call with the arguments the
harness actually sent, and the agent's folder (inbox, calendar, files, memory)
updating as it changes.

    python -m webui.server            then open http://127.0.0.1:8765

Binds loopback only. One run at a time, in a subprocess (webui/runner.py), so
Stop always works and the process-global harness switches can't collide.
"""
import http.server
import json
import mimetypes
import os
import queue
import shutil
import socket
import subprocess
import sys
import threading
import time
import urllib.parse
import webbrowser

import requests

HERE = os.path.dirname(os.path.abspath(__file__))
PROJECT = os.path.dirname(HERE)
STATIC = os.path.join(HERE, "static")
AGENTS_DIR = os.path.join(PROJECT, "agents")
OLLAMA_URL = "http://127.0.0.1:11434"


def default_model():
    """The model the setup screen checks for when the page names none: the
    first agent's configured tag."""
    for name in agent_folders():
        try:
            with open(os.path.join(AGENTS_DIR, name, "config.json"),
                      encoding="utf-8") as f:
                return json.load(f).get("model") or ""
        except Exception:
            continue
    return ""
DEFAULT_PORT = 8765

sys.path.insert(0, PROJECT)
from harness import chat  # noqa: E402
from harness import mcp_config  # noqa: E402
from harness import profiles  # noqa: E402
from webui import preflight  # noqa: E402

# Rough per-size guidance for the picker; the machine, not the harness, decides.
# Model facts from openrouter.ai, generated into model_catalog.json rather than
# fetched. The product's whole claim is that nothing leaves the machine, so the
# UI must not reach the network to describe a model. Regenerate with
# tools/refresh_catalog.py when the shipped model list changes.
def _load_catalog():
    path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "model_catalog.json")
    try:
        with open(path, encoding="utf-8") as f:
            return json.load(f).get("models", {})
    except (OSError, ValueError):
        return {}


CATALOG = _load_catalog()


def catalog_for(tag):
    """Exact tag first, then the bare family, so llama3.1 matches llama3.1:8b."""
    if tag in CATALOG:
        return CATALOG[tag]
    base = tag.split(":")[0]
    for key, val in CATALOG.items():
        if key.split(":")[0] == base:
            return val
    return {}


SPEED_HINT = {
    "1b": ("instant", "Fast enough to feel live. Makes the most mistakes — the best "
                      "place to watch the harness repair a call."),
    "3b": ("quick", "A few seconds per step. The sweet spot for demos."),
    "8b": ("steady", "Tens of seconds per step. Noticeably more reliable."),
    "14b": ("slow", "Heavy on CPU-only machines. Strong tool discipline."),
    "32b": ("very slow", "Minutes per step. Fits in 32 GB RAM, tests your patience."),
}

PRESET_TASKS = [
    "Summarize my Wednesday meetings and message Jordan with the list",
    "Find a free hour on Thursday and book it as Deep work",
    "Turn Dana's Q3 sales numbers into a PowerPoint deck",
    "Build a spreadsheet of my July receipts with a total",
    "Reply to Mia about the Northwind kickoff and add it to my calendar",
    "Remember that I prefer meetings after 14:00 and never on Fridays",
]


# ----------------------------------------------------------------- agents ----

def agent_folders():
    if not os.path.isdir(AGENTS_DIR):
        return []
    out = []
    for name in sorted(os.listdir(AGENTS_DIR), key=lambda n: (len(n), n)):
        if name.startswith("_") or name.startswith("."):
            continue
        cfg_path = os.path.join(AGENTS_DIR, name, "config.json")
        if os.path.isfile(cfg_path):
            out.append(name)
    return out


def read_config(agent):
    with open(os.path.join(AGENTS_DIR, agent, "config.json"), encoding="utf-8-sig") as f:
        return json.load(f)


def agent_dir(agent):
    """Resolve an agent id to its folder, refusing anything else."""
    if agent not in agent_folders():
        raise ValueError(f"unknown agent {agent!r}")
    return os.path.join(AGENTS_DIR, agent)


def installed_tags():
    try:
        r = requests.get(f"{OLLAMA_URL}/api/tags", timeout=3)
        r.raise_for_status()
        return {m["name"]: m.get("size", 0) for m in r.json().get("models", [])}
    except Exception:
        return None  # None = server unreachable, {} = up with no models


def tag_installed(tag, tags):
    """Ollama treats llama3.1:8b and llama3.1:latest as different tags even when
    they share blobs, so only an exact match counts — except that a bare name
    means :latest."""
    if not tags:
        return False
    return tag in tags or (":" not in tag and f"{tag}:latest" in tags)


def agent_list():
    tags = installed_tags()
    out = []
    for name in agent_folders():
        cfg = read_config(name)
        folder = os.path.join(AGENTS_DIR, name)
        files_dir = os.path.join(folder, "workspace", "files")
        logs_dir = os.path.join(folder, "logs")
        mem_path = os.path.join(folder, "memory", "memory.jsonl")
        speed, blurb = SPEED_HINT.get(name, ("", ""))
        profile = profiles.for_model(cfg["model"], cfg.get("harness"))
        out.append({
            "id": name,
            "name": cfg.get("name", name),
            "model": cfg["model"],
            "note": cfg.get("note", ""),
            "speed": speed,
            "blurb": blurb,
            "catalog": catalog_for(cfg["model"]),
            "profile": profile.to_dict(),
            "installed": tag_installed(cfg["model"], tags),
            "files": len(os.listdir(files_dir)) if os.path.isdir(files_dir) else 0,
            "runs": len([f for f in os.listdir(logs_dir) if f.startswith("run_")])
                    if os.path.isdir(logs_dir) else 0,
            "memories": sum(1 for _ in open(mem_path, encoding="utf-8"))
                        if os.path.isfile(mem_path) else 0,
        })
    # Models the catalog knows about that are not installed. The rail offers
    # them for download, so a machine with one model is not a dead end.
    have = {a["model"] for a in out}
    available = [dict(v, tag=k) for k, v in CATALOG.items()
                 if k not in have and not tag_installed(k, tags)]
    available.sort(key=lambda m: m["tag"])
    return {"agents": out, "ollama": tags is not None, "presets": PRESET_TASKS,
            "available": available,
            "project": PROJECT, "installed_models": sorted(tags or {})}


# -------------------------------------------------------------- workspace ----

def workspace(agent):
    """The agent's folder as the browser shows it — same shape the runner emits
    during a run, so the panel renders identically live and at rest."""
    folder = agent_dir(agent)
    state_path = os.path.join(folder, "workspace", "state.json")
    files_dir = os.path.join(folder, "workspace", "files")
    mem_path = os.path.join(folder, "memory", "memory.jsonl")
    logs_dir = os.path.join(folder, "logs")

    state = {}
    if os.path.isfile(state_path):
        try:
            with open(state_path, encoding="utf-8") as f:
                state = json.load(f)
        except ValueError:
            state = {}
    if not state:  # never run: show the fixtures it will start from
        from harness.world import CALENDAR, EMAILS
        state = {"emails": [dict(e) for e in EMAILS], "events": [dict(e) for e in CALENDAR]}

    files = []
    if os.path.isdir(files_dir):
        for name in sorted(os.listdir(files_dir)):
            path = os.path.join(files_dir, name)
            if os.path.isfile(path):
                st = os.stat(path)
                files.append({"name": name, "size": st.st_size, "mtime": st.st_mtime})

    memory = []
    if os.path.isfile(mem_path):
        with open(mem_path, encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if line:
                    try:
                        memory.append(json.loads(line)["fact"])
                    except (ValueError, KeyError):
                        pass

    logs = []
    if os.path.isdir(logs_dir):
        for name in sorted(os.listdir(logs_dir), reverse=True):
            if name.startswith("run_") and name.endswith(".json"):
                logs.append({"name": name,
                             "mtime": os.path.getmtime(os.path.join(logs_dir, name))})

    return {
        # Newest first, the way every mail client and the agent's own
        # list_emails tool order them. This panel rendered them in insertion
        # order, so a newly arrived email appeared at the BOTTOM of the inbox,
        # under nine older ones.
        "emails": sorted(state.get("emails", []), key=lambda e: e.get("date", ""),
                         reverse=True),
        "sent": state.get("sent_emails", []),
        "events": sorted(state.get("events", []), key=lambda e: (e["date"], e["start"])),
        "messages": state.get("messages", []),
        "reminders": state.get("reminders", []),
        "files": files,
        "memory": memory,
        "logs": logs[:25],
        "folder": folder,
    }


def workspace_file(agent, name):
    files_dir = os.path.join(agent_dir(agent), "workspace", "files")
    path = os.path.abspath(os.path.join(files_dir, os.path.basename(str(name))))
    if os.path.dirname(path) != os.path.abspath(files_dir) or not os.path.isfile(path):
        raise ValueError(f"no such file {name!r}")
    return path


EMU_PER_PT = 12700
MAX_SHEET_ROWS = 300
MAX_SHEET_COLS = 40


def _cell_text(val, numeric, cell):
    """What the cell says, honouring its number format.

    The preview printed str(value), so a figure written with a #,##0 format sat
    in the browser as 1240000 while Excel showed 1,240,000 for the same file.
    On a deck built from an email that quoted "$1,240,000" it read as though the
    agent had mangled the number. Only the grouped formats this app writes are
    handled; anything else is left alone rather than guessed at.
    """
    if val is None:
        return ""
    if not numeric:
        return str(val)
    fmt = (cell.number_format or "") if cell is not None else ""
    if "#,##0" in fmt:
        return f"{val:,.2f}" if ".00" in fmt else f"{val:,.0f}"
    return str(val)


def _solid_fill(shape):
    """'#RRGGBB' for a solid-filled autoshape, else None. Accent bars and rules
    carry the deck's design; without this they render as invisible boxes."""
    try:
        color = shape.fill.fore_color
        if color.type is not None and color.rgb is not None:
            return "#" + str(color.rgb)
    except Exception:
        pass
    return None


def _run_color(para):
    try:
        for r in para.runs:
            if r.font.color is not None and r.font.color.rgb is not None:
                return "#" + str(r.font.color.rgb)
        if para.font.color is not None and para.font.color.rgb is not None:
            return "#" + str(para.font.color.rgb)
    except Exception:
        pass
    return None


def _pptx_preview(path):
    """Slide geometry, not a text dump — every shape's real position, size and
    font, as fractions of the slide, so the browser can lay the deck out the way
    PowerPoint would. Fractions (not absolute units) keep it resolution-free:
    the front end scales one number and everything follows."""
    import base64

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    sw = prs.slide_width or 9144000
    sh_ = prs.slide_height or 6858000
    slides = []
    for slide in prs.slides:
        shapes = []
        for shape in slide.shapes:
            item = {"x": (shape.left or 0) / sw, "y": (shape.top or 0) / sh_,
                    "w": (shape.width or 0) / sw, "h": (shape.height or 0) / sh_}
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = shape.image
                    item.update(kind="picture", src="data:%s;base64,%s" % (
                        img.content_type, base64.b64encode(img.blob).decode()))
                    shapes.append(item)
                    continue
            except Exception:
                pass
            if getattr(shape, "has_table", False):
                item.update(kind="table",
                            rows=[[c.text for c in row.cells] for row in shape.table.rows])
                shapes.append(item)
                continue
            if not getattr(shape, "has_text_frame", False):
                item.update(kind="box", fill=_solid_fill(shape))
                shapes.append(item)
                continue
            try:
                is_title = slide.shapes.title is not None and shape == slide.shapes.title
            except Exception:
                is_title = False
            paras = []
            for p in shape.text_frame.paragraphs:
                text = "".join(r.text for r in p.runs) or p.text or ""
                if not text.strip():
                    continue
                # run-level first, then the paragraph default the runs inherit,
                # and only then a guess
                size = next((r.font.size.pt for r in p.runs if r.font.size is not None), None)
                if size is None and p.font.size is not None:
                    size = p.font.size.pt
                bold = next((r.font.bold for r in p.runs if r.font.bold is not None), None)
                if bold is None:
                    bold = p.font.bold
                color = _run_color(p)
                paras.append({
                    "text": text,
                    "level": p.level or 0,
                    # python-pptx leaves size/bold None when inherited from the
                    # layout; fall back to sane title/body defaults rather than
                    # rendering everything at one size.
                    "size": size if size else (30.0 if is_title else 18.0),
                    "bold": bool(bold) if bold is not None else is_title,
                    "align": getattr(p.alignment, "name", "").lower() or None,
                    "color": color,
                })
            # An autoshape has a text frame even when empty, so a solid accent
            # bar lands here rather than in the no-text-frame branch above.
            if paras:
                item.update(kind="text", title=is_title, paragraphs=paras)
            else:
                item.update(kind="box", fill=_solid_fill(shape))
            shapes.append(item)
        slides.append({"shapes": shapes})
    return {"kind": "pptx", "name": os.path.basename(path),
            "w_pt": sw / EMU_PER_PT, "h_pt": sh_ / EMU_PER_PT, "slides": slides}


def _xlsx_preview(path):
    """A real grid: column letters, row numbers, widths, merges, bold, and both
    sides of a formula cell (the formula and its cached result)."""
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    wb_f = load_workbook(path, data_only=False)   # formulas as written
    wb_v = load_workbook(path, data_only=True)    # values as last computed
    sheets = []
    for ws in wb_f.worksheets:
        wsv = wb_v[ws.title]
        n_rows = min(ws.max_row or 1, MAX_SHEET_ROWS)
        n_cols = min(ws.max_column or 1, MAX_SHEET_COLS)
        widths = []
        for i in range(1, n_cols + 1):
            dim = ws.column_dimensions.get(get_column_letter(i))
            widths.append(round((dim.width if dim and dim.width else 9.5) * 7.5))
        merges = [{"r": m.min_row, "c": m.min_col,
                   "rs": m.max_row - m.min_row + 1, "cs": m.max_col - m.min_col + 1}
                  for m in ws.merged_cells.ranges
                  if m.min_row <= n_rows and m.min_col <= n_cols]
        rows = []
        for r in range(1, n_rows + 1):
            cells = []
            for c in range(1, n_cols + 1):
                cf, cv = ws.cell(row=r, column=c), wsv.cell(row=r, column=c)
                raw = cf.value
                formula = raw if isinstance(raw, str) and raw.startswith("=") else None
                val = cv.value
                if val is None and formula is None:
                    val = raw
                # openpyxl does not evaluate formulas, and a file Excel has
                # never opened carries no cached result — so show the formula
                # itself rather than an empty cell.
                if val is None and formula is not None:
                    val = formula
                numeric = isinstance(val, (int, float)) and not isinstance(val, bool)
                align = (cf.alignment.horizontal if cf.alignment else None) or \
                        ("right" if numeric else None)
                cells.append({"v": _cell_text(val, numeric, cf), "f": formula,
                              "b": bool(cf.font and cf.font.bold), "a": align})
            rows.append(cells)
        sheets.append({"sheet": ws.title, "rows": rows, "widths": widths,
                       "merges": merges, "cols": n_cols,
                       "truncated": (ws.max_row or 0) > n_rows})
    return {"kind": "xlsx", "name": os.path.basename(path), "sheets": sheets}


def preview(agent, name):
    """Render a generated file in the browser instead of making the user open
    PowerPoint — the whole point is to see what the agent produced."""
    path = workspace_file(agent, name)
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pptx":
        return _pptx_preview(path)
    if ext == ".xlsx":
        return _xlsx_preview(path)
    with open(path, "rb") as f:
        blob = f.read(20000)
    if b"\x00" in blob[:2000]:
        return {"kind": "binary", "name": os.path.basename(path),
                "size": os.path.getsize(path)}
    return {"kind": "text", "name": os.path.basename(path),
            "text": blob.decode("utf-8", errors="replace")}


def ensure_idle(agent):
    """Refuse to touch an agent's folder while its run is live.

    Reset mid-run did not just race - it silently UNDID itself: the running
    subprocess holds the world in memory and snapshots on exit, so the freshly
    deleted state.json reappeared seconds later carrying the pre-reset world.
    The user saw the reset succeed and then saw the old inbox come back.
    Same contract as starting a second run: one thing owns the folder at a
    time, and the answer is 409, not corruption."""
    run = RUNS.current
    if run and run.agent == agent and run.proc.poll() is None:
        raise RuntimeError(f"{agent} is mid-run - stop the run first, then reset")


def reset_agent(agent, what):
    folder = agent_dir(agent)
    done = []
    targets = {
        "world": os.path.join(folder, "workspace", "state.json"),
        "memory": os.path.join(folder, "memory", "memory.jsonl"),
    }
    for key, path in targets.items():
        if key in what and os.path.isfile(path):
            os.remove(path)
            done.append(key)
    for key, path in (("files", os.path.join(folder, "workspace", "files")),
                      ("logs", os.path.join(folder, "logs"))):
        if key in what and os.path.isdir(path):
            shutil.rmtree(path)
            os.makedirs(path, exist_ok=True)
            done.append(key)
    return done


def reveal(path):
    if sys.platform == "darwin":
        subprocess.Popen(["open", path])
    elif sys.platform.startswith("win"):
        os.startfile(path)  # noqa: S606 - local dev convenience
    else:
        subprocess.Popen(["xdg-open", path])


# -------------------------------------------------------------------- runs ----

class Run:
    """One agent subprocess, its event log, and everyone watching it."""

    def __init__(self, rid, agent, task, proc, options):
        self.id = rid
        self.agent = agent
        self.task = task
        self.proc = proc
        self.options = options
        self.started = time.time()
        self.events = []
        self.subs = []
        self.status = "running"
        self.lock = threading.Lock()

    def add(self, event):
        with self.lock:
            self.events.append(event)
            item = (len(self.events) - 1, event)
            subs = list(self.subs)
        for q in subs:
            q.put(item)

    def subscribe(self, after=-1):
        """Register a watcher and hand back everything it has not seen. `after`
        comes from Last-Event-ID, so a reconnect resumes instead of replaying."""
        q = queue.Queue()
        with self.lock:
            backlog = list(enumerate(self.events))[after + 1:]
            self.subs.append(q)
        return q, backlog

    def unsubscribe(self, q):
        with self.lock:
            if q in self.subs:
                self.subs.remove(q)

    def answer(self, cid, allow):
        try:
            self.proc.stdin.write(json.dumps({"id": cid, "allow": bool(allow)}) + "\n")
            self.proc.stdin.flush()
            return True
        except (OSError, ValueError):
            return False

    def stop(self):
        if self.proc.poll() is None:
            self.proc.terminate()
            self.status = "stopped"


class Runs:
    def __init__(self):
        self.current = None
        self.next_id = 1
        self.lock = threading.Lock()

    def start(self, agent, task, options):
        with self.lock:
            if self.current and self.current.proc.poll() is None:
                raise RuntimeError(f"{self.current.agent} is already running — "
                                   "stop it first (one model at a time).")
            cmd = [sys.executable, "-u", "-m", "webui.runner",
                   "--agent", agent, "--task", task]
            if options.get("root"):
                cmd += ["--root", options["root"]]
            if options.get("shell"):
                cmd.append("--shell")
            if options.get("yolo"):
                cmd.append("--yolo")
            if options.get("with_office"):
                cmd.append("--with-office")
            if options.get("tiers"):
                cmd.append("--tiers")
            if options.get("small"):
                cmd += ["--small", options["small"]]
            if options.get("deep"):
                cmd += ["--deep", options["deep"]]
            if options.get("max_calls"):
                cmd += ["--max-calls", str(int(options["max_calls"]))]
            if options.get("model"):
                cmd += ["--model", options["model"]]
            if options.get("thread"):
                cmd += ["--thread", options["thread"]]
            if options.get("mcp"):
                cmd += ["--mcp", options["mcp"]]
            if options.get("mcp_mode"):
                cmd += ["--mcp-mode", options["mcp_mode"]]
            env = dict(os.environ, PYTHONUNBUFFERED="1", PYTHONIOENCODING="utf-8")
            proc = subprocess.Popen(cmd, cwd=PROJECT, env=env, text=True,
                                    encoding="utf-8", errors="replace", bufsize=1,
                                    stdin=subprocess.PIPE, stdout=subprocess.PIPE,
                                    stderr=subprocess.PIPE)
            if options.get("thread"):
                chat.append(agent_dir(agent), options["thread"], "user", task)
            run = Run(self.next_id, agent, task, proc, options)
            self.next_id += 1
            self.current = run
        threading.Thread(target=self._pump, args=(run,), daemon=True).start()
        return run

    def _pump(self, run):
        stderr = []
        threading.Thread(target=lambda: stderr.extend(run.proc.stderr),
                         daemon=True).start()
        for line in run.proc.stdout:
            line = line.strip()
            if not line:
                continue
            try:
                run.add(json.loads(line))
            except ValueError:
                run.add({"t": "stdout", "text": line})
        code = run.proc.wait()
        time.sleep(0.05)  # let the stderr reader drain
        if run.status == "running":
            run.status = "finished" if code == 0 else "failed"
        if code not in (0, -15) and run.status != "stopped":
            tail = "".join(stderr)[-1500:].strip()
            run.add({"t": "error", "message": f"the run exited with code {code}",
                     "trace": tail})
        # The agent's reply to the conversation is its done() summary. Recorded
        # here rather than in the runner so a crashed or stopped run still
        # leaves an honest turn in the thread instead of a silent gap.
        tid = run.options.get("thread")
        if tid:
            end = next((e for e in reversed(run.events) if e.get("t") == "end"), None)
            spoken = chat.said(run.events)
            if spoken:
                # A run that already spoke keeps what it said. Appending the
                # done summary underneath would paste a summary of the message
                # below the message.
                reply = "\n\n".join(spoken)
            elif end and end.get("summary"):
                reply = end["summary"]
            elif run.status == "stopped":
                reply = "(stopped)"
            else:
                reply = "(the run ended without a summary)"
            chat.append(agent_dir(run.agent), tid, "assistant", reply, run=run.id)
        run.add({"t": "closed", "status": run.status, "code": code})


RUNS = Runs()


# ------------------------------------------------------------------ server ----

class Handler(http.server.BaseHTTPRequestHandler):
    protocol_version = "HTTP/1.1"
    server_version = "AgentLab"

    def log_message(self, *args):
        pass  # the console belongs to the run banner, not to request noise

    # ---- helpers ----
    def send_json(self, obj, status=200):
        body = json.dumps(obj, ensure_ascii=False, default=str).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.send_header("Cache-Control", "no-store")
        self.end_headers()
        self.wfile.write(body)

    def send_bytes(self, blob, ctype, extra=()):
        self.send_response(200)
        self.send_header("Content-Type", ctype)
        self.send_header("Content-Length", str(len(blob)))
        self.send_header("Cache-Control", "no-store")
        for k, v in extra:
            self.send_header(k, v)
        self.end_headers()
        self.wfile.write(blob)

    def body_json(self):
        length = int(self.headers.get("Content-Length") or 0)
        if not length:
            return {}
        return json.loads(self.rfile.read(length).decode("utf-8"))

    def query(self):
        parts = urllib.parse.urlparse(self.path)
        return parts.path, {k: v[0] for k, v in urllib.parse.parse_qs(parts.query).items()}

    # ---- GET ----
    def do_GET(self):
        path, q = self.query()
        try:
            if path in ("/", "/index.html"):
                return self.static_file("index.html")
            if path.startswith("/static/"):
                return self.static_file(path[len("/static/"):])
            # Both must be served from the root, not /static/: a service worker
            # can only control a scope at or below its own path, and the install
            # prompt needs the manifest's scope to cover start_url.
            if path == "/sw.js":
                return self.static_file("sw.js")
            if path == "/manifest.webmanifest":
                return self.static_file("manifest.webmanifest")
            if path == "/api/agents":
                return self.send_json(agent_list())
            if path == "/api/workspace":
                return self.send_json(workspace(q.get("agent", "")))
            if path == "/api/preview":
                return self.send_json(preview(q.get("agent", ""), q.get("name", "")))
            if path == "/api/download":
                fpath = workspace_file(q.get("agent", ""), q.get("name", ""))
                ctype = mimetypes.guess_type(fpath)[0] or "application/octet-stream"
                with open(fpath, "rb") as f:
                    blob = f.read()
                return self.send_bytes(blob, ctype, extra=[
                    ("Content-Disposition",
                     f'attachment; filename="{os.path.basename(fpath)}"')])
            if path == "/api/log":
                folder = agent_dir(q.get("agent", ""))
                name = os.path.basename(q.get("name", ""))
                with open(os.path.join(folder, "logs", name), encoding="utf-8") as f:
                    return self.send_json(json.load(f))
            if path == "/api/threads":
                return self.send_json({"threads": chat.threads(agent_dir(q.get("agent", "")))})
            if path == "/api/thread":
                folder = agent_dir(q.get("agent", ""))
                return self.send_json({"id": q.get("id", ""),
                                       "messages": chat.messages(folder, q.get("id", ""))})
            if path == "/api/mcp":
                # The registry, for the run panel's account picker. Setup notes
                # come along so the UI can say what a server needs before it works.
                return self.send_json([
                    {"name": name, "summary": summary,
                     "setup": mcp_config.setup_notes(name)}
                    for name, summary in mcp_config.available()])
            if path == "/api/setup":
                # What a first run needs before it can work, plus any download
                # in flight. Polled by the setup screen.
                model = q.get("model") or default_model()
                checks = preflight.check(model)
                return self.send_json({"checks": checks,
                                       "ready": preflight.ready(checks),
                                       "pull": preflight.PULL})
            if path == "/api/status":
                run = RUNS.current
                return self.send_json({"run": run.id if run else None,
                                       "agent": run.agent if run else None,
                                       "status": run.status if run else "idle"})
            if path == "/api/events":
                return self.stream_events(q)
        except ValueError as e:
            return self.send_json({"error": str(e)}, 400)
        except FileNotFoundError as e:
            return self.send_json({"error": str(e)}, 404)
        except Exception as e:  # a broken panel shouldn't take the server down
            return self.send_json({"error": f"{type(e).__name__}: {e}"}, 500)
        self.send_json({"error": "not found"}, 404)

    # ---- POST ----
    def do_POST(self):
        path, _ = self.query()
        try:
            body = self.body_json()
            if path == "/api/thread/new":
                folder = agent_dir(body.get("agent", ""))
                return self.send_json({"id": chat.create(folder, body.get("task", ""))})
            if path == "/api/thread/delete":
                folder = agent_dir(body.get("agent", ""))
                return self.send_json({"deleted": chat.delete(folder, body.get("id", ""))})
            if path == "/api/run":
                agent = body.get("agent", "")
                task = (body.get("task") or "").strip()
                agent_dir(agent)  # validates
                if not task:
                    raise ValueError("give the agent a task first")
                root = (body.get("root") or "").strip()
                if root and not os.path.isdir(os.path.expanduser(root)):
                    raise ValueError(f"working folder {root} does not exist")
                options = {"root": os.path.expanduser(root) if root else None,
                           "shell": body.get("shell"), "yolo": body.get("yolo"),
                           "with_office": body.get("with_office"),
                           "tiers": body.get("tiers"), "small": body.get("small"),
                           "deep": body.get("deep"), "max_calls": body.get("max_calls"),
                           "model": (body.get("model") or "").strip() or None,
                           "thread": (body.get("thread") or "").strip() or None,
                           "mcp": ",".join(body.get("mcp") or []) or None,
                           "mcp_mode": body.get("mcp_mode") or None}
                run = RUNS.start(agent, task, options)
                return self.send_json({"run": run.id, "agent": agent})
            if path == "/api/setup/fix":
                # Only the named, safe actions in preflight.FIXES: install this
                # interpreter's packages, start Ollama, pull a model. Installing
                # Ollama itself stays a link the user clicks.
                result = preflight.apply_fix(body.get("action", ""),
                                             {"tag": body.get("tag") or default_model()})
                return self.send_json({"ok": True, "result": result})
            if path == "/api/stop":
                run = RUNS.current
                if run:
                    run.stop()
                return self.send_json({"ok": True})
            if path == "/api/confirm":
                run = RUNS.current
                ok = bool(run) and run.answer(int(body.get("id", 0)), body.get("allow"))
                return self.send_json({"ok": ok})
            if path == "/api/reset":
                what = set(body.get("what") or [])
                ensure_idle(body.get("agent", ""))
                return self.send_json({"cleared": reset_agent(body.get("agent", ""), what)})
            if path == "/api/reveal":
                agent = body.get("agent", "")
                sub = body.get("sub") or ""
                target = os.path.join(agent_dir(agent), *sub.split("/")) if sub \
                    else agent_dir(agent)
                if not os.path.exists(target):
                    raise ValueError("that folder does not exist yet")
                reveal(target)
                return self.send_json({"ok": True, "path": target})
        except RuntimeError as e:
            return self.send_json({"error": str(e)}, 409)
        except ValueError as e:
            return self.send_json({"error": str(e)}, 400)
        except Exception as e:
            return self.send_json({"error": f"{type(e).__name__}: {e}"}, 500)
        self.send_json({"error": "not found"}, 404)

    # ---- static ----
    def static_file(self, rel):
        path = os.path.abspath(os.path.join(STATIC, rel))
        if not path.startswith(os.path.abspath(STATIC)) or not os.path.isfile(path):
            return self.send_json({"error": "not found"}, 404)
        with open(path, "rb") as f:
            blob = f.read()
        ctype = mimetypes.guess_type(path)[0] or "text/plain"
        self.send_bytes(blob, f"{ctype}; charset=utf-8" if "text" in ctype
                        or "javascript" in ctype else ctype)

    # ---- SSE ----
    def open_stream(self):
        self.send_response(200)
        self.send_header("Content-Type", "text/event-stream; charset=utf-8")
        self.send_header("Cache-Control", "no-store")
        self.send_header("Connection", "close")
        self.end_headers()

    def push(self, obj, index=None):
        prefix = f"id: {index}\n" if index is not None else ""
        self.wfile.write(f"{prefix}data: {json.dumps(obj, ensure_ascii=False, default=str)}\n\n"
                         .encode("utf-8"))
        self.wfile.flush()

    def stream_events(self, q):
        run = RUNS.current
        want = q.get("run")
        if not run or (want and str(run.id) != str(want)):
            self.open_stream()
            self.push({"t": "closed", "status": "gone"})
            return
        try:
            after = int(self.headers.get("Last-Event-ID"))
        except (TypeError, ValueError):
            after = -1
        sub, backlog = run.subscribe(after)
        self.open_stream()
        self.close_connection = True
        try:
            for index, event in backlog:
                self.push(event, index)
            while True:
                try:
                    index, event = sub.get(timeout=10)
                except queue.Empty:
                    self.wfile.write(b": keepalive\n\n")
                    self.wfile.flush()
                    continue
                self.push(event, index)
                if event.get("t") == "closed":
                    break
        except (BrokenPipeError, ConnectionResetError):
            pass
        finally:
            run.unsubscribe(sub)

class Server(http.server.ThreadingHTTPServer):
    daemon_threads = True
    allow_reuse_address = True


def free_port(start=DEFAULT_PORT):
    for port in range(start, start + 20):
        with socket.socket() as s:
            if s.connect_ex(("127.0.0.1", port)) != 0:
                return port
    return start


def main():
    sys.path.insert(0, PROJECT)
    port = free_port(int(os.environ.get("AGENT_LAB_PORT", DEFAULT_PORT)))
    url = f"http://127.0.0.1:{port}"
    tags = installed_tags()
    print(f"\n  Agent Lab  →  {url}")
    print(f"  project    {PROJECT}")
    if tags is None:
        print("  ollama     NOT RUNNING — start it, then reload the page")
    else:
        print(f"  ollama     up, {len(tags)} model(s) installed")
    print("  agents     " + ", ".join(agent_folders()))
    print("\n  Ctrl-C to stop.\n")
    if os.environ.get("AGENT_LAB_NO_BROWSER") != "1":
        threading.Timer(0.6, lambda: webbrowser.open(url)).start()
    try:
        Server(("127.0.0.1", port), Handler).serve_forever()
    except KeyboardInterrupt:
        if RUNS.current:
            RUNS.current.stop()
        print("\n  stopped.\n")


if __name__ == "__main__":
    main()
